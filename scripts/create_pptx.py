#!/usr/bin/env python3
"""
Create PowerPoint presentation for Poker Dream Ecosystem
Based on the PDF: Poker_Dream_The_Definitive_Ecosystem.pdf
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml

# Color scheme matching the PDF
COLORS = {
    'gold': RgbColor(0xC9, 0xA2, 0x27),  # Gold/Mustard
    'green': RgbColor(0x2D, 0x5A, 0x3D),  # Dark Green
    'light_green': RgbColor(0x4A, 0x7C, 0x59),
    'black': RgbColor(0x1A, 0x1A, 0x1A),
    'dark_gray': RgbColor(0x33, 0x33, 0x33),
    'light_gray': RgbColor(0xF5, 0xF5, 0xF0),  # Off-white background
    'white': RgbColor(0xFF, 0xFF, 0xFF),
    'red': RgbColor(0xC0, 0x39, 0x2B),
}

def set_slide_background(slide, color):
    """Set solid background color for a slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_text(slide, text, left, top, width, height, font_size=44, bold=True, color=COLORS['black'], align=PP_ALIGN.LEFT):
    """Add a text box with title styling"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Arial"
    p.alignment = align
    return txBox

def add_body_text(slide, text, left, top, width, height, font_size=18, color=COLORS['dark_gray'], align=PP_ALIGN.LEFT, bold=False):
    """Add a text box with body styling"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.name = "Arial"
    p.font.bold = bold
    p.alignment = align
    return txBox

def add_bullet_points(slide, items, left, top, width, height, font_size=14, color=COLORS['dark_gray']):
    """Add bullet point list"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Arial"
        p.space_after = Pt(8)
    return txBox

def add_card(slide, title, body_items, left, top, width, height, title_color=COLORS['gold']):
    """Add a card-style content block"""
    # Card background
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['white']
    shape.line.color.rgb = RgbColor(0xDD, 0xDD, 0xDD)

    # Title
    add_body_text(slide, title, left + Inches(0.2), top + Inches(0.15),
                  width - Inches(0.4), Inches(0.4), font_size=16, color=title_color, bold=True)

    # Body bullets
    if body_items:
        add_bullet_points(slide, body_items, left + Inches(0.2), top + Inches(0.5),
                         width - Inches(0.4), height - Inches(0.6), font_size=12)

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
    prs.slide_height = Inches(7.5)

    # Use blank layout
    blank_layout = prs.slide_layouts[6]

    # =========================================================================
    # SLIDE 1: Title Slide
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1, COLORS['light_gray'])

    add_title_text(slide1, "Poker Dream: The Definitive\nHome for Serious Poker",
                   Inches(0.8), Inches(1.5), Inches(11), Inches(2), font_size=48, align=PP_ALIGN.CENTER)

    add_body_text(slide1, "An ecosystem designed to empower individual mastery\nand foster collective community.",
                  Inches(0.8), Inches(3.5), Inches(11), Inches(1), font_size=20, align=PP_ALIGN.CENTER)

    add_body_text(slide1, "Poker Dream", Inches(5.5), Inches(6.5), Inches(2.5), Inches(0.5),
                  font_size=24, bold=True, align=PP_ALIGN.CENTER)

    # =========================================================================
    # SLIDE 2: Problem Statement
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2, COLORS['light_gray'])

    add_title_text(slide2, "The Modern Poker Player is Alone and Flying Blind",
                   Inches(0.5), Inches(0.4), Inches(12), Inches(1), font_size=36)

    add_body_text(slide2, "The serious player's journey is fragmented and filled with guesswork. They operate without the professional tools or centralized community that exists in every other serious pursuit.",
                  Inches(0.5), Inches(1.2), Inches(12), Inches(0.8), font_size=16)

    # Three problem cards
    problems = [
        ("Inconsistent Data & Guesswork",
         ["Players struggle with tedious, error-prone tracking",
          "Memory bias leads to underestimating losses",
          "They don't know their true ROI or hourly rate"],
         '"I use Excel but it\'s a pain. I just want to tap a few buttons and be done."'),
        ("Information Overload",
         ["Poker news, schedules, and streams are scattered",
          "Content spread across operator sites, YouTube, Twitch",
          "Leads to missed content and information lag"],
         '"I have to switch between 5+ apps just to follow a single tournament series."'),
        ("A Disconnected Community",
         ["Community is fractured across generic social platforms",
          "Lacking a dedicated space for focused discussion",
          "Hard to find like-minded poker friends"],
         '"General social media doesn\'t get it. It\'s hard to find like-minded poker friends."')
    ]

    for i, (title, points, quote) in enumerate(problems):
        left = Inches(0.5 + i * 4.1)
        add_body_text(slide2, title, left, Inches(2.3), Inches(3.8), Inches(0.6),
                      font_size=18, bold=True, color=COLORS['black'])
        add_bullet_points(slide2, points, left, Inches(2.9), Inches(3.8), Inches(2.5), font_size=13)
        add_body_text(slide2, quote, left, Inches(5.5), Inches(3.8), Inches(1),
                      font_size=12, color=COLORS['gold'])

    # =========================================================================
    # SLIDE 3: Two-Pillar Strategy
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3, COLORS['light_gray'])

    add_title_text(slide3, "Our Two-Pillar Strategy to Own the Player Lifecycle",
                   Inches(0.5), Inches(0.4), Inches(12), Inches(1), font_size=36)

    add_body_text(slide3, "We will capture the entire player journey by addressing their two fundamental needs in sequence: first, the need for self-mastery, and second, the need for community connection.",
                  Inches(0.5), Inches(1.2), Inches(12), Inches(0.8), font_size=16)

    add_body_text(slide3, "The Player", Inches(5.8), Inches(2.2), Inches(2), Inches(0.4),
                  font_size=16, align=PP_ALIGN.CENTER)

    # Mastery box
    shape1 = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(3), Inches(5), Inches(3))
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = COLORS['white']
    shape1.line.color.rgb = COLORS['gold']

    add_body_text(slide3, "Mastery", Inches(1.2), Inches(3.3), Inches(4.6), Inches(0.5),
                  font_size=24, bold=True, color=COLORS['gold'], align=PP_ALIGN.CENTER)
    add_body_text(slide3, "Empowering players with professional-grade tools to turn their hobby into a data-driven business.",
                  Inches(1.2), Inches(3.9), Inches(4.6), Inches(1), font_size=14, align=PP_ALIGN.CENTER)
    add_body_text(slide3, "Poker Dream Bankroll", Inches(1.5), Inches(5.3), Inches(4), Inches(0.4),
                  font_size=14, bold=True, color=COLORS['gold'], align=PP_ALIGN.CENTER)

    # Community box
    shape2 = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(3), Inches(5), Inches(3))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = COLORS['white']
    shape2.line.color.rgb = COLORS['green']

    add_body_text(slide3, "Community", Inches(7.2), Inches(3.3), Inches(4.6), Inches(0.5),
                  font_size=24, bold=True, color=COLORS['green'], align=PP_ALIGN.CENTER)
    add_body_text(slide3, "Connecting empowered players to a vibrant, centralized platform for news, events, and interaction.",
                  Inches(7.2), Inches(3.9), Inches(4.6), Inches(1), font_size=14, align=PP_ALIGN.CENTER)
    add_body_text(slide3, "Poker Dream Community", Inches(7.5), Inches(5.3), Inches(4), Inches(0.4),
                  font_size=14, bold=True, color=COLORS['green'], align=PP_ALIGN.CENTER)

    # =========================================================================
    # SLIDE 4: Pillar 1 - Mastery
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4, COLORS['light_gray'])

    add_body_text(slide4, "PILLAR 1: MASTERY", Inches(0.3), Inches(2.5), Inches(0.5), Inches(3),
                  font_size=12, color=COLORS['gold'], bold=True)

    add_title_text(slide4, "Pillar 1: Mastery with\nPoker Dream Bankroll",
                   Inches(5), Inches(0.8), Inches(7.5), Inches(1.5), font_size=40, color=COLORS['gold'])

    add_body_text(slide4, "Track. Analyze. Improve.", Inches(5), Inches(2.3), Inches(7), Inches(0.5),
                  font_size=20, color=COLORS['dark_gray'])

    add_body_text(slide4, 'Meet "Serious Sam."', Inches(5), Inches(3.2), Inches(7), Inches(0.5),
                  font_size=22, bold=True)

    add_body_text(slide4, "He's a skilled player, but he struggles to prove his hobby is profitable. He uses tedious spreadsheets and wants to know if he's truly ready to move up in stakes. He needs to eliminate the guesswork.",
                  Inches(5), Inches(3.8), Inches(7.5), Inches(1.2), font_size=16)

    # Highlight box
    shape = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5), Inches(5.2), Inches(7.5), Inches(1.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RgbColor(0xFF, 0xF8, 0xE7)
    shape.line.color.rgb = COLORS['gold']

    add_body_text(slide4, "For serious players like Sam, Poker Dream Bankroll provides comprehensive session logging, bankroll management, and advanced analytics to manage their poker finances like a business.",
                  Inches(5.2), Inches(5.4), Inches(7.1), Inches(1.2), font_size=15)

    # =========================================================================
    # SLIDE 5: Bankroll Features
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5, COLORS['light_gray'])

    add_title_text(slide5, "Turning Poker From a Gamble Into a Business",
                   Inches(0.5), Inches(0.4), Inches(12), Inches(0.8), font_size=36)

    features = [
        ("Eliminate Guesswork", COLORS['gold'], [
            "Fast, Mobile-First Session Logging: Log cash games and tournaments in under 90 seconds",
            "Dual Game Support: Track both cash game metrics (hourly rate, win rate) and tournament metrics (ROI, ITM%)",
            "Comprehensive Data Capture: Record venue, stakes, buy-ins, duration, position, and payouts"
        ]),
        ("Manage Your Finances Professionally", COLORS['gold'], [
            "Dedicated Bankroll Management: Track deposits, withdrawals, and session results separately",
            "Visualize Your Growth: See long-term trends with interactive bankroll charts",
            "Professional Tax Reporting: Generate clean CSV and PDF exports for tax filing"
        ]),
        ("Unlock True Performance", COLORS['gold'], [
            "Advanced Analytics: Instantly see your true ROI, hourly rate, and win rate",
            "Filter & Compare: Identify your most profitable venues, stakes, and game types",
            "Built-in Pro Tools: Integrated ICM, Equity, and Hand Range calculators"
        ])
    ]

    for i, (title, color, points) in enumerate(features):
        left = Inches(0.5 + i * 4.1)
        add_body_text(slide5, title, left, Inches(1.5), Inches(3.8), Inches(0.6),
                      font_size=20, bold=True, color=color)
        add_bullet_points(slide5, points, left, Inches(2.2), Inches(3.8), Inches(4.5), font_size=13)

    # =========================================================================
    # SLIDE 6: Pillar 2 - Community
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6, COLORS['light_gray'])

    add_body_text(slide6, "PILLAR 2: COMMUNITY", Inches(12.5), Inches(2.5), Inches(0.6), Inches(3),
                  font_size=12, color=COLORS['green'], bold=True)

    add_title_text(slide6, "Once Empowered,\nPlayers Seek\nConnection.",
                   Inches(0.8), Inches(0.8), Inches(6), Inches(2.5), font_size=44)

    add_body_text(slide6, 'Introducing Pillar 2: The Poker Dream Community Platform.\n"Your Poker Community."',
                  Inches(0.8), Inches(3.2), Inches(6), Inches(1), font_size=18)

    add_body_text(slide6, 'This is "Tournament Tommy."', Inches(0.8), Inches(4.5), Inches(6), Inches(0.5),
                  font_size=20, bold=True)

    add_body_text(slide6, "He uses the Bankroll app and knows his stats. Now, he wants to follow the World Series of Poker in real-time.\n\nHe's frustrated jumping between Twitter, a streaming site, and a clunky forum to discuss hands and track results.\n\nHe's missing the biggest moments.",
                  Inches(0.8), Inches(5.1), Inches(6), Inches(2), font_size=15)

    # =========================================================================
    # SLIDE 7: Community Features
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7, COLORS['light_gray'])

    add_title_text(slide7, "The Centralized Hub for the Entire Poker World",
                   Inches(0.5), Inches(0.4), Inches(12), Inches(0.8), font_size=36)

    comm_features = [
        ("The Single Source of Truth", COLORS['green'], [
            "Aggregated News Feed: Curated news from top sources, all in one place",
            "Comprehensive Events Calendar: Track major tournament series like WSOP, WPT, and EPT",
            "Live Stream Discovery: Find every major poker stream from YouTube, Twitch, and Facebook"
        ]),
        ("The Global Poker Table", COLORS['green'], [
            "Dedicated Social Feed: A poker-focused feed to share hands, stories, and content",
            "Real-Time Chat Rooms: Join public rooms, create private groups, and send direct messages",
            "Event-Specific Lounges: Discuss hands and follow action in dedicated chat rooms"
        ]),
        ("Never Miss a Moment", COLORS['green'], [
            "Personalized Notifications: Get alerts for breaking news, event reminders, and player updates",
            "Follow System: Connect with and follow other players, pros, and content creators",
            "Video Highlights: Catch up quickly with a curated feed of key hands and interviews"
        ])
    ]

    for i, (title, color, points) in enumerate(comm_features):
        left = Inches(0.5 + i * 4.1)
        add_body_text(slide7, title, left, Inches(1.5), Inches(3.8), Inches(0.6),
                      font_size=20, bold=True, color=color)
        add_bullet_points(slide7, points, left, Inches(2.2), Inches(3.8), Inches(4.5), font_size=13)

    # =========================================================================
    # SLIDE 8: Ecosystem Flywheel
    # =========================================================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8, COLORS['light_gray'])

    add_title_text(slide8, "The Poker Dream Ecosystem Flywheel",
                   Inches(0.5), Inches(0.3), Inches(12), Inches(0.8), font_size=36, align=PP_ALIGN.CENTER)

    flywheel_steps = [
        ("1. ACQUIRE & EMPOWER", Inches(8), Inches(1.2), COLORS['gold'],
         "Players are drawn to the Poker Dream Bankroll App to track performance and improve their game.\n\nThis is their entry point to self-mastery."),
        ("2. ENGAGE & CONNECT", Inches(8), Inches(3.8), COLORS['green'],
         "Empowered, data-savvy players graduate to the Poker Dream Community App to discuss strategy, follow events, and connect with peers.\n\nTheir individual success fuels a desire for collective engagement."),
        ("3. GROW THE NETWORK", Inches(0.5), Inches(3.8), COLORS['green'],
         "The vibrant Community attracts more players through network effects. New members are then introduced to the Bankroll App.\n\nThe community becomes a powerful, organic acquisition channel."),
        ("4. CREATE A MOAT", Inches(0.5), Inches(1.2), COLORS['gold'],
         "This self-reinforcing loop increases user retention and lifetime value (LTV), creating a deep, defensible moat that is difficult for competitors to replicate.\n\nWe own the entire player journey.")
    ]

    for title, left, top, color, desc in flywheel_steps:
        add_body_text(slide8, title, left, top, Inches(4.5), Inches(0.4), font_size=16, bold=True, color=color)
        add_body_text(slide8, desc, left, top + Inches(0.4), Inches(4.5), Inches(2), font_size=12)

    # =========================================================================
    # SLIDE 9: App Preview
    # =========================================================================
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide9, COLORS['light_gray'])

    add_title_text(slide9, "A Glimpse into the Poker Dream Suite",
                   Inches(0.5), Inches(0.3), Inches(12), Inches(0.8), font_size=36, align=PP_ALIGN.CENTER)

    add_body_text(slide9, "Poker Dream Bankroll", Inches(1.5), Inches(1.3), Inches(4), Inches(0.5),
                  font_size=24, bold=True, color=COLORS['gold'], align=PP_ALIGN.CENTER)

    add_body_text(slide9, "Poker Dream Community", Inches(7.5), Inches(1.3), Inches(4), Inches(0.5),
                  font_size=24, bold=True, color=COLORS['green'], align=PP_ALIGN.CENTER)

    # Placeholder for app screenshots
    add_body_text(slide9, "[Dashboard, Session Logging, Analytics screens]",
                  Inches(1), Inches(2.5), Inches(5), Inches(4), font_size=14, align=PP_ALIGN.CENTER)
    add_body_text(slide9, "[Home, Events Calendar, Chat screens]",
                  Inches(7), Inches(2.5), Inches(5), Inches(4), font_size=14, align=PP_ALIGN.CENTER)

    # =========================================================================
    # SLIDE 10: Market Opportunity
    # =========================================================================
    slide10 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide10, COLORS['light_gray'])

    add_title_text(slide10, "Targeting a Valuable and Underserved Market",
                   Inches(0.5), Inches(0.3), Inches(12), Inches(0.8), font_size=36)

    # Market stats
    add_body_text(slide10, "Market Opportunity", Inches(0.5), Inches(1.2), Inches(5.5), Inches(0.4),
                  font_size=18, bold=True)
    add_body_text(slide10, "$6.9B+", Inches(0.5), Inches(1.7), Inches(2.5), Inches(0.7),
                  font_size=48, bold=True, color=COLORS['green'])
    add_body_text(slide10, "Global Online Poker Market (2024)", Inches(0.5), Inches(2.4), Inches(2.5), Inches(0.4),
                  font_size=12)
    add_body_text(slide10, "100M+", Inches(3.2), Inches(1.7), Inches(2.5), Inches(0.7),
                  font_size=48, bold=True, color=COLORS['green'])
    add_body_text(slide10, "Poker Players Worldwide", Inches(3.2), Inches(2.4), Inches(2.5), Inches(0.4),
                  font_size=12)

    add_body_text(slide10, "Existing tools are fragmented, desktop-focused, or have a dated mobile UX. There is no dominant, mobile-first ecosystem.",
                  Inches(0.5), Inches(3.1), Inches(5.5), Inches(0.8), font_size=14)

    # Primary Audience
    add_body_text(slide10, "Primary Target Audience", Inches(7), Inches(1.2), Inches(5.5), Inches(0.4),
                  font_size=18, bold=True, align=PP_ALIGN.CENTER)
    add_body_text(slide10, '"Serious Sam" & "Pro Paula"', Inches(7), Inches(1.7), Inches(5.5), Inches(0.4),
                  font_size=14, bold=True, align=PP_ALIGN.CENTER)
    add_body_text(slide10, "The recreational grinders and semi-professionals who treat poker as a business or a profitable hobby.\n\nAlready tracking performance (or know they should be), tech-savvy, and willing to pay for tools that provide a clear ROI.",
                  Inches(7), Inches(2.2), Inches(5.5), Inches(1.5), font_size=13, align=PP_ALIGN.CENTER)

    # Secondary Audience
    add_body_text(slide10, "Secondary Audience", Inches(0.5), Inches(4.2), Inches(5.5), Inches(0.4),
                  font_size=18, bold=True)
    add_body_text(slide10, '"Tournament Tommy" & "Social Sarah"', Inches(0.5), Inches(4.7), Inches(5.5), Inches(0.4),
                  font_size=14, bold=True)
    add_body_text(slide10, "The avid fans and community connectors who drive network effects.\n\nHighly engaged with poker content and community, driving daily active use and attracting new users to the ecosystem.",
                  Inches(0.5), Inches(5.2), Inches(5.5), Inches(1.5), font_size=13)

    # =========================================================================
    # SLIDE 11: Monetization
    # =========================================================================
    slide11 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide11, COLORS['light_gray'])

    add_title_text(slide11, "A Dual Monetization Strategy",
                   Inches(0.5), Inches(0.3), Inches(12), Inches(0.8), font_size=36, align=PP_ALIGN.CENTER)

    # Bankroll monetization
    add_body_text(slide11, "Poker Dream Bankroll:\nPremium Tools (SaaS)", Inches(0.5), Inches(1.2), Inches(5.5), Inches(0.8),
                  font_size=22, bold=True, color=COLORS['gold'], align=PP_ALIGN.CENTER)
    add_body_text(slide11, "Model: Freemium\nFree Plan: Core tracking features with a 3-month session history limit.",
                  Inches(0.5), Inches(2.1), Inches(5.5), Inches(0.8), font_size=14, align=PP_ALIGN.CENTER)

    # Bankroll Pro box
    shape = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3), Inches(5), Inches(2.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RgbColor(0xFF, 0xF8, 0xE7)
    shape.line.color.rgb = COLORS['gold']

    add_body_text(slide11, "Pro Plan: $4.99/month or $49.99 Lifetime. Unlocks:",
                  Inches(1), Inches(3.15), Inches(4.6), Inches(0.4), font_size=13, bold=True)
    add_bullet_points(slide11, ["Unlimited Session History", "Advanced Statistics & Filtering",
                                "PDF & CSV Data Export", "Integrated Poker Tools (ICM, Equity)", "Ad-Free Experience"],
                      Inches(1), Inches(3.5), Inches(4.6), Inches(1.8), font_size=12)

    add_body_text(slide11, "Goal: 15% free-to-paid conversion rate.",
                  Inches(0.5), Inches(5.5), Inches(5.5), Inches(0.4), font_size=14, bold=True, align=PP_ALIGN.CENTER)

    # Community monetization
    add_body_text(slide11, "Poker Dream Community:\nEngagement & Scale (Media)", Inches(7), Inches(1.2), Inches(5.5), Inches(0.8),
                  font_size=22, bold=True, color=COLORS['green'], align=PP_ALIGN.CENTER)
    add_body_text(slide11, "Model: Ad-Supported with Premium Option\nFree Plan: Full access to all community features, supported by non-intrusive ads.",
                  Inches(7), Inches(2.1), Inches(5.5), Inches(0.8), font_size=14, align=PP_ALIGN.CENTER)

    # Community Premium box
    shape2 = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.3), Inches(3), Inches(5), Inches(2.2))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = RgbColor(0xE8, 0xF5, 0xE9)
    shape2.line.color.rgb = COLORS['green']

    add_body_text(slide11, "Premium Plan: $8.99/month. Provides:",
                  Inches(7.5), Inches(3.15), Inches(4.6), Inches(0.4), font_size=13, bold=True)
    add_bullet_points(slide11, ["Completely Ad-Free Experience", "Exclusive Content & AMA Sessions",
                                "Premium Profile Badge", "Early Access to New Features"],
                      Inches(7.5), Inches(3.5), Inches(4.6), Inches(1.8), font_size=12)

    add_body_text(slide11, "Goal: 5% of MAU converting to premium.",
                  Inches(7), Inches(5.5), Inches(5.5), Inches(0.4), font_size=14, bold=True, align=PP_ALIGN.CENTER)

    # =========================================================================
    # SLIDE 12: 6-Month Goals
    # =========================================================================
    slide12 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide12, COLORS['light_gray'])

    add_title_text(slide12, "Measuring Our Ascent: Key 6-Month Goals",
                   Inches(0.5), Inches(0.3), Inches(12), Inches(0.8), font_size=36)

    add_body_text(slide12, "Our success is defined by a clear set of metrics focused on adoption, engagement, and sustainable monetization.",
                  Inches(0.5), Inches(1.0), Inches(12), Inches(0.5), font_size=16)

    metrics = [
        ("Adoption & Reach", [
            "Bankroll Downloads: 50,000+",
            "Community Downloads: 150,000+",
            "App Store Rating: 4.5+ Stars"
        ]),
        ("User Engagement", [
            "Bankroll: 70% of users log at least 1 session/week",
            "Community: 27%+ DAU/MAU Ratio",
            "Retention: Achieve 40% D30 retention"
        ]),
        ("Monetization", [
            "Bankroll Conversion: 15% Free-to-Pro conversion",
            "ARPU (Ecosystem): Target blended ARPU of $2.50+"
        ]),
        ("Product Quality", [
            "Crash-Free Rate: 99.5%+",
            "API Uptime: 99.9%",
            "NPS Score: 45+"
        ])
    ]

    for i, (title, items) in enumerate(metrics):
        row = i // 2
        col = i % 2
        left = Inches(0.5 + col * 6.3)
        top = Inches(1.8 + row * 2.5)

        add_body_text(slide12, title, left, top, Inches(5.8), Inches(0.5), font_size=20, bold=True)
        for j, item in enumerate(items):
            parts = item.split(": ")
            if len(parts) == 2:
                add_body_text(slide12, parts[0] + ":", left, top + Inches(0.5 + j * 0.4), Inches(3.5), Inches(0.4), font_size=14)
                add_body_text(slide12, parts[1], left + Inches(3.5), top + Inches(0.5 + j * 0.4), Inches(2), Inches(0.4),
                              font_size=14, bold=True, color=COLORS['green'])

    # =========================================================================
    # SLIDE 13: Phased Rollout
    # =========================================================================
    slide13 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide13, COLORS['light_gray'])

    add_title_text(slide13, "The Path to the Ecosystem: A Phased Rollout",
                   Inches(0.5), Inches(0.3), Inches(12), Inches(0.8), font_size=36)

    phases = [
        ("Phase 1: Foundation", "Months 1-2", COLORS['gold'],
         "Focus: Launch Poker Dream Bankroll MVP.",
         "Deliverables: Core session logging, bankroll tracking, and basic statistics.",
         "Goal: Validate the core tool and build a user base of serious players."),
        ("Phase 2: Enhancement", "Months 3-4", RgbColor(0x8B, 0xC3, 0x4A),
         "Focus: Add Bankroll Pro features.",
         "Deliverables: Advanced analytics, data export (CSV/PDF), and integrated Poker Tools.",
         "Goal: Drive monetization and solidify Bankroll as the market-leading tool."),
        ("Phase 3: Connection", "Months 5-6", COLORS['green'],
         "Focus: Launch Poker Dream Community MVP.",
         "Deliverables: Home feed (news/video), events calendar, live streams, and social feed with chat.",
         "Goal: Capture the network effect and begin building the flywheel."),
        ("Phase 4: Integration & Growth", "Months 7-12", RgbColor(0x60, 0x7D, 0x8B),
         "Focus: Deepen the ecosystem.",
         "Deliverables: Cross-promotion between apps, unified user profiles, and advanced community features.",
         "Goal: Accelerate the flywheel and establish market leadership.")
    ]

    for i, (title, timeline, color, focus, deliverables, goal) in enumerate(phases):
        left = Inches(0.3 + i * 3.2)

        # Phase box
        shape = slide13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.3), Inches(3), Inches(5.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS['white']
        shape.line.color.rgb = color

        add_body_text(slide13, title, left + Inches(0.15), Inches(1.45), Inches(2.7), Inches(0.4),
                      font_size=16, bold=True)
        add_body_text(slide13, timeline, left + Inches(0.15), Inches(1.9), Inches(2.7), Inches(0.3),
                      font_size=12, color=color)
        add_body_text(slide13, focus, left + Inches(0.15), Inches(2.4), Inches(2.7), Inches(0.6), font_size=11, bold=True)
        add_body_text(slide13, deliverables, left + Inches(0.15), Inches(3.1), Inches(2.7), Inches(1.5), font_size=11)
        add_body_text(slide13, goal, left + Inches(0.15), Inches(5), Inches(2.7), Inches(1.2), font_size=11, bold=True)

    # =========================================================================
    # SLIDE 14: Why We Will Win
    # =========================================================================
    slide14 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide14, COLORS['light_gray'])

    add_title_text(slide14, "Why We Will Win",
                   Inches(0.5), Inches(0.3), Inches(12), Inches(0.8), font_size=40)

    add_body_text(slide14, "While other trackers exist, they are single-purpose tools, not ecosystems. Our competition is fragmented, mobile-second, and lacks a unified vision for the modern player.",
                  Inches(0.5), Inches(1.0), Inches(12), Inches(0.7), font_size=16)

    # Comparison table
    headers = ["Feature", "Typical Competitor", "Poker Dream"]
    rows = [
        ["Strategy", "Single App (Tracking Tool)", "Integrated Ecosystem (Tool + Platform)"],
        ["User Experience", "Dated UI, Cluttered Navigation", "Mobile-First, Premium Editorial Design"],
        ["Value Prop", "Basic Tracking", "Tracking + Pro Tools + Community + Content"],
        ["Community", "None / External Forums", "Built-in, Real-Time, Event-Driven"],
        ["Content", "None", "Aggregated News, Streams & Live Results"]
    ]

    # Table headers
    for j, header in enumerate(headers):
        left = Inches(0.5 + j * 4)
        width = Inches(3.8) if j > 0 else Inches(2)
        add_body_text(slide14, header, left, Inches(2), width, Inches(0.4), font_size=14, bold=True)

    # Table rows
    for i, row in enumerate(rows):
        for j, cell in enumerate(row):
            left = Inches(0.5 + j * 4)
            width = Inches(3.8) if j > 0 else Inches(2)
            color = COLORS['green'] if j == 2 else COLORS['dark_gray']
            add_body_text(slide14, cell, left, Inches(2.5 + i * 0.6), width, Inches(0.5), font_size=13, color=color)

    # Bottom highlight
    shape = slide14.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.8), Inches(12), Inches(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RgbColor(0xFF, 0xF8, 0xE7)
    shape.line.color.rgb = COLORS['green']

    add_body_text(slide14, "Our two-pillar strategy creates a powerful flywheel, building a defensible moat that a simple tracking tool cannot compete with.",
                  Inches(0.7), Inches(6), Inches(11.6), Inches(0.7), font_size=16, bold=True)

    # =========================================================================
    # SLIDE 15: Closing
    # =========================================================================
    slide15 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide15, COLORS['black'])

    add_title_text(slide15, "We Are Not Just Building Apps.\nWe Are Building the Home for Poker.",
                   Inches(0.5), Inches(2), Inches(12), Inches(2), font_size=40, color=COLORS['white'], align=PP_ALIGN.CENTER)

    add_body_text(slide15, "Poker Dream will become the essential, indispensable ecosystem for the modern poker player. We will empower their individual journey to mastery and connect them to the global community, creating a single, definitive destination for the sport we love.",
                  Inches(1.5), Inches(4.5), Inches(10), Inches(1.5), font_size=16, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Save the presentation
    output_path = "/Users/clifflai/Downloads/Poker_Dream_Ecosystem.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()
